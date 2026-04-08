"""
MindMap Sales Collateral — RAG Ingestion Pipeline
==================================================
Single script: extract → chunk → enrich → embed → upsert to Qdrant

Usage:
    python ingest.py                    # full run
    python ingest.py --dry-run          # extract + chunk only, no API calls
    python ingest.py --recreate         # drop and recreate Qdrant collection first
    python ingest.py --file "foo.pdf"   # process a single file (useful for testing)
"""

from __future__ import annotations

import sys
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

import argparse
import base64
import hashlib
import json
import re
import time
import uuid
import warnings
from pathlib import Path
from typing import Dict, List, Set

import fitz          # PyMuPDF — PDF→image render for VLM pages
import pypdf
from docx import Document
from openai import OpenAI
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance,
    PayloadSchemaType,
    PointStruct,
    VectorParams,
)

from config import (
    CHUNK_OVERLAP_WORDS,
    CHUNK_SIZE_WORDS,
    CLIENT_NAME_MAP,
    COLLECTION_NAME,
    DATA_DIR,
    EMBEDDING_BATCH,
    EMBEDDING_DIM,
    EMBEDDING_MODEL,
    IMAGE_DPI,
    MIN_CHUNK_WORDS,
    ONEPAGER_THRESHOLD,
    OPENAI_API_KEY,
    QDRANT_API_KEY,
    QDRANT_URL,
    SKIP_DUPLICATES,
    SKIP_EXTENSIONS,
    SKIP_FOLDERS,
    SKIP_PDF_USE_PPTX,
    UPSERT_BATCH,
    VLM_MODEL,
    VLM_REQUIRED_FILES,
    FOLDER_TO_DOCTYPE,
    FOLDER_TO_VERTICAL,
)

warnings.filterwarnings("ignore")

# ─────────────────────────────────────────────────────────────────────────────
# Lazy singletons
# ─────────────────────────────────────────────────────────────────────────────

_openai: OpenAI | None = None
_qdrant: QdrantClient | None = None


def openai_client() -> OpenAI:
    global _openai
    if _openai is None:
        _openai = OpenAI(api_key=OPENAI_API_KEY)
    return _openai


def qdrant_client() -> QdrantClient:
    global _qdrant
    if _qdrant is None:
        kwargs: dict = {"url": QDRANT_URL}
        if QDRANT_API_KEY:
            kwargs["api_key"] = QDRANT_API_KEY
        _qdrant = QdrantClient(**kwargs)
    return _qdrant


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 1 — EXTRACTION
# Three paths: pypdf (text PDFs) · pptx · vlm (image-only pages/docs)
# ═════════════════════════════════════════════════════════════════════════════

def extract_file(path: Path) -> List[Dict]:
    """
    Route a file to the correct extractor.
    Returns list of: { page: int, text: str, extraction_method: str }
    """
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".docx":
        return _extract_docx(path)
    return []


# ── PDF ───────────────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> List[Dict]:
    # Files confirmed fully image-rendered → skip pypdf entirely
    if path.name in VLM_REQUIRED_FILES:
        return _all_pages_vlm(path)

    pages = _pypdf_extract(path)

    # Per-page fallback: if pypdf got < MIN_CHUNK_WORDS on a page, use VLM
    image_pages = [p for p in pages if len(p["text"].split()) < MIN_CHUNK_WORDS]
    if image_pages:
        pages = _patch_with_vlm(path, pages)

    return pages


def _pypdf_extract(path: Path) -> List[Dict]:
    results = []
    try:
        reader = pypdf.PdfReader(str(path))
        for i, page in enumerate(reader.pages):
            raw  = page.extract_text() or ""
            text = raw.encode("utf-8", errors="replace").decode("utf-8").strip()
            results.append({"page": i + 1, "text": text, "extraction_method": "pypdf"})
    except Exception as e:
        print(f"    [WARN] pypdf error on {path.name}: {e}")
    return results


def _patch_with_vlm(path: Path, pages: List[Dict]) -> List[Dict]:
    """Re-extract image-sparse pages via GPT-4o, leave text-rich pages alone."""
    try:
        doc = fitz.open(str(path))
        for page_data in pages:
            if len(page_data["text"].split()) < MIN_CHUNK_WORDS:
                fitz_page = doc[page_data["page"] - 1]
                page_data["text"]             = _vlm_single_page(fitz_page)
                page_data["extraction_method"] = "vlm"
        doc.close()
    except Exception as e:
        print(f"    [WARN] VLM patch failed for {path.name}: {e}")
    return pages


def _all_pages_vlm(path: Path) -> List[Dict]:
    """Send every page through GPT-4o vision (for fully visual PDFs)."""
    results = []
    try:
        doc = fitz.open(str(path))
        for i, page in enumerate(doc):
            results.append({
                "page": i + 1,
                "text": _vlm_single_page(page),
                "extraction_method": "vlm",
            })
        doc.close()
    except Exception as e:
        print(f"    [WARN] Full-VLM extraction failed for {path.name}: {e}")
    return results


def _vlm_single_page(page: fitz.Page) -> str:
    """Render one PDF page → PNG → GPT-4o → extracted text."""
    mat = fitz.Matrix(IMAGE_DPI / 72, IMAGE_DPI / 72)
    pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
    b64 = base64.b64encode(pix.tobytes("png")).decode()

    resp = openai_client().chat.completions.create(
        model=VLM_MODEL,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": (
                        "Extract all text, numbers, headings, bullet points, and key metrics "
                        "from this sales collateral image exactly as written. "
                        "Preserve logical reading order. Output only the extracted content."
                    ),
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"},
                },
            ],
        }],
        max_tokens=1000,
    )
    return resp.choices[0].message.content.strip()


# ── PPTX ──────────────────────────────────────────────────────────────────

def _extract_pptx(path: Path) -> List[Dict]:
    results = []
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                # Regular text frames
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                # Table shapes — iterate cells individually
                if shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            texts.append(" | ".join(row_cells))
            text = "\n".join(texts).strip()
            if text:
                results.append({"page": i + 1, "text": text, "extraction_method": "pptx"})
    except Exception as e:
        print(f"    [WARN] PPTX extraction failed for {path.name}: {e}")
    return results


# ── DOCX ──────────────────────────────────────────────────────────────────

def _extract_docx(path: Path) -> List[Dict]:
    try:
        doc   = Document(str(path))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        paras.append(t)
        full = "\n".join(paras)
        if full.strip():
            return [{"page": 1, "text": full, "extraction_method": "docx"}]
    except Exception as e:
        print(f"    [WARN] DOCX extraction failed: {e}")
    return []


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 2 — SMART CHUNKING
# Strategy selected per doc_type:
#   case_study  → section-based (Problem/Solution/Benefits headers)
#   short doc   → single chunk
#   everything  → page-merge + sliding window
# ═════════════════════════════════════════════════════════════════════════════

_SECTION_HEADERS = [
    r"problem\s+statement", r"our\s+approach", r"approach",
    r"our\s+solution",      r"solution",       r"benefits?",
    r"outcomes?",           r"results?",        r"key\s+results?",
    r"roi\b",               r"challenges?",     r"project\s+description",
    r"overview",            r"organization\s+profiled",
    r"business\s+processes?\s+impacted",        r"objectives?",
    r"case\s+study\s+brief",                    r"what\s+we\s+automated",
    r"the\s+method",        r"impact\s+and\s+reach",
    r"stages?\s+of\s+transformation",
]
_SECTION_RE = re.compile(
    r"(?:^|\n)(" + "|".join(_SECTION_HEADERS) + r")\s*[:\-]?\s*\n",
    re.IGNORECASE | re.MULTILINE,
)


def chunk_pages(pages: List[Dict], doc_type: str) -> List[Dict]:
    pages = [p for p in pages if len(p.get("text", "").split()) >= MIN_CHUNK_WORDS]
    if not pages:
        return []

    full_text   = "\n\n".join(p["text"] for p in pages)
    total_words = len(full_text.split())

    # ── Single chunk (one-pagers / very short) ────────────────────────────
    if total_words <= ONEPAGER_THRESHOLD:
        return [{
            "text":             full_text,
            "section_type":     "full_document",
            "page_numbers":     [p["page"] for p in pages],
            "extraction_method": pages[0].get("extraction_method", "pypdf"),
        }]

    # ── Section-based (case studies) ─────────────────────────────────────
    if doc_type == "case_study":
        section_chunks = _section_split(full_text, pages)
        if section_chunks:
            return section_chunks
        # No recognisable headers — fall through to sliding window

    # ── Page-merge + sliding window (everything else) ─────────────────────
    return _page_merge_chunks(pages)


def _section_split(full_text: str, pages: List[Dict]) -> List[Dict]:
    matches = list(_SECTION_RE.finditer(full_text))
    if len(matches) < 2:
        return []

    all_page_nums = [p["page"] for p in pages]
    method        = pages[0].get("extraction_method", "pypdf")
    chunks        = []

    # Capture intro text BEFORE the first section header (title, client context, etc.)
    intro = full_text[:matches[0].start()].strip()
    if len(intro.split()) >= MIN_CHUNK_WORDS:
        chunks.append({
            "text":             intro,
            "section_type":     "introduction",
            "page_numbers":     all_page_nums,
            "extraction_method": method,
        })

    for i, match in enumerate(matches):
        header = match.group(1).strip().lower()
        start  = match.end()
        end    = matches[i + 1].start() if i + 1 < len(matches) else len(full_text)
        body   = full_text[start:end].strip()

        if len(body.split()) < MIN_CHUNK_WORDS:
            continue

        if len(body.split()) > CHUNK_SIZE_WORDS * 1.5:
            chunks.extend(_sliding_window(body, header, all_page_nums, method))
        else:
            chunks.append({
                "text":             body,
                "section_type":     header,
                "page_numbers":     all_page_nums,
                "extraction_method": method,
            })

    return chunks


def _page_merge_chunks(pages: List[Dict]) -> List[Dict]:
    """Merge consecutive pages up to word budget, then sliding-window each buffer."""
    chunks: List[Dict] = []
    buf_text   = ""
    buf_pages: List[int] = []
    buf_method = pages[0].get("extraction_method", "pypdf")

    for page_data in pages:
        text   = page_data["text"].strip()
        method = page_data.get("extraction_method", "pypdf")
        pg     = page_data["page"]

        combined_wc = len((buf_text + " " + text).split()) if buf_text else len(text.split())

        if combined_wc <= CHUNK_SIZE_WORDS:
            buf_text   = (buf_text + "\n\n" + text).strip() if buf_text else text
            buf_pages.append(pg)
            buf_method = method
        else:
            if len(buf_text.split()) >= MIN_CHUNK_WORDS:
                chunks.extend(_sliding_window(buf_text, "content", list(buf_pages), buf_method))
            buf_text   = text
            buf_pages  = [pg]
            buf_method = method

    if len(buf_text.split()) >= MIN_CHUNK_WORDS:
        chunks.extend(_sliding_window(buf_text, "content", list(buf_pages), buf_method))

    return chunks


def _sliding_window(text: str, section_type: str, page_numbers: List[int], method: str) -> List[Dict]:
    words  = text.split()
    chunks = []
    start  = 0
    while start < len(words):
        end        = min(start + CHUNK_SIZE_WORDS, len(words))
        chunk_text = " ".join(words[start:end])
        if len(chunk_text.split()) >= MIN_CHUNK_WORDS:
            chunks.append({
                "text":             chunk_text,
                "section_type":     section_type,
                "page_numbers":     page_numbers,
                "extraction_method": method,
            })
        if end == len(words):
            break
        start += CHUNK_SIZE_WORDS - CHUNK_OVERLAP_WORDS
    return chunks


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 3 — METADATA ENRICHMENT
# Auto-detect doc_type, industry vertical, function area, solution type, ROI flag
# ═════════════════════════════════════════════════════════════════════════════

_VERTICAL_KEYWORDS: Dict[str, List[str]] = {
    "BFSI":          ["bank", "banking", "insurance", "fintech", "mortgage",
                      "trade processing", "neo bank", "capital markets", "wealth management"],
    "Healthcare":    ["hospital", "patient", "healthcare", "clinical", "medical",
                      "physician", "caregiver", "ehr", "emr", "diagnosis", "home care"],
    "Pharma":        ["pharma", "pharmaceutical", "drug manufacturing", "intas", "piramal"],
    "FA":            ["fp&a", "accounts payable", "accounts receivable", "invoice",
                      "journal entries", "reconciliation", "finance", "accounting", "budget"],
    "HR":            ["recruitment", "onboarding", "payroll", "human resources",
                      "employee", "caregiver onboarding"],
    "IT":            ["itsm", "it helpdesk", "service desk", "incident management",
                      "it operations", "servicenow"],
    "Manufacturing": ["manufacturing", "quality control", "production", "plant", "batch"],
    "SCM":           ["supply chain", "demand planning", "procurement", "inventory", "order management"],
    "Telecom":       ["telecom", "telecommunications"],
    "Logistics":     ["logistics", "transportation", "freight", "shipping"],
    "Government":    ["government", "public sector"],
    "Retail":        ["retail"],
    "Education":     ["higher education", "university", "elearning", "training program"],
    "SAP":           ["sap", "s/4hana"],
    "Aviation":      ["airline", "airport", "aviation"],
}

_FUNCTION_KEYWORDS: Dict[str, List[str]] = {
    "AP":             ["accounts payable", "invoice processing", "vendor payment", "ap automation"],
    "AR":             ["accounts receivable", "cash receipts", "collections", "billing reconciliation", "cash applications"],
    "FPA":            ["fp&a", "financial planning", "reporting automation", "budget", "mis reporting"],
    "HR":             ["recruitment", "onboarding", "payroll", "hr digital assistant"],
    "ITSM":           ["itsm", "service desk", "incident", "it helpdesk", "change request"],
    "Trade":          ["trade processing", "trade settlement", "trade capture"],
    "SCM":            ["demand planning", "supply chain", "order management", "vendor reconciliation"],
    "Customer Service": ["customer service", "contact center", "customer onboarding", "loyalty"],
    "Compliance":     ["regulatory", "compliance", "kyc", "aml", "pre-authorization"],
    "Clinical":       ["diagnosis", "patient access", "patient admission", "appointment scheduling"],
    "GL/Accounting":  ["journal entries", "accruals", "general ledger", "book closure", "gl clearance"],
}

_ROI_RE = re.compile(
    r"(\d+\s*%\s*(reduction|increase|improvement|saving|faster|accuracy|fte))"
    r"|(\$[\d,]+)"
    r"|(roi\s+of\s+\d+)"
    r"|(payback[:\s]+\d+)"
    r"|(annual\s+(savings?|benefit))"
    r"|(fte\s+reduc)",
    re.IGNORECASE,
)


def enrich_chunk(path: Path, chunk: Dict, doc_first500: str) -> Dict:
    """Attach metadata to a chunk dict."""
    combined   = chunk["text"] + " " + doc_first500
    text_lower = combined.lower()
    folder_parts = [p.lower() for p in path.parts]

    # doc_type
    doc_type = "capability_deck"
    for folder_key, dtype in FOLDER_TO_DOCTYPE.items():
        if folder_key.lower() in " ".join(folder_parts):
            doc_type = dtype
            break

    # industry verticals — from folder path first, then text
    verticals: Set[str] = set()
    for part in folder_parts:
        for fkey, vlist in FOLDER_TO_VERTICAL.items():
            if fkey.lower() in part:
                verticals.update(vlist)
    for vertical, keywords in _VERTICAL_KEYWORDS.items():
        if any(kw in text_lower for kw in keywords):
            verticals.add(vertical)
    if not verticals:
        verticals.add("General")

    # function areas
    functions = [
        func for func, keywords in _FUNCTION_KEYWORDS.items()
        if any(kw in text_lower for kw in keywords)
    ]

    # solution types
    solutions: List[str] = []
    if any(k in text_lower for k in ["rpa", "robotic process", "automation bot", " bot "]):
        solutions.append("RPA")
    if any(k in text_lower for k in ["machine learning", " ai ", "nlp", "ocr", "icr", "cognitive", "artificial intelligence"]):
        solutions.append("AI/ML")
    if any(k in text_lower for k in ["analytics", "dashboard", "business intelligence", "data lake", "power bi", "tableau", "qlikview"]):
        solutions.append("Analytics")
    if any(k in text_lower for k in ["aws", "gcp", "azure", "cloud"]):
        solutions.append("Cloud")
    if any(k in text_lower for k in ["chatbot", "virtual assistant", "conversational ai"]):
        solutions.append("Chatbot")
    if any(k in text_lower for k in ["staff augmentation", "resource fulfilment", "talent"]):
        solutions.append("Staff Augmentation")
    if not solutions:
        solutions.append("RPA")

    # client name — match against file name + folder path
    searchable = path.name.lower() + " " + " ".join(folder_parts)
    client_name = None
    for key, name in CLIENT_NAME_MAP.items():
        if key in searchable:
            client_name = name
            break

    return {
        **chunk,
        "doc_type":          doc_type,
        "industry_vertical": list(verticals),
        "function_area":     functions,
        "solution_type":     solutions,
        "has_roi_metrics":   bool(_ROI_RE.search(chunk["text"])),
        "client_name":       client_name,
        "file_name":         path.name,
        "file_path":         _relative_path(path),
        "source_folder":     _source_folder(path),
    }


def _relative_path(path: Path) -> str:
    try:
        return str(path.relative_to(DATA_DIR))
    except ValueError:
        return path.name


def _source_folder(path: Path) -> str:
    """Immediate subfolder of Sales Collateral."""
    try:
        parts = path.relative_to(DATA_DIR).parts
        return parts[0] if len(parts) > 1 else "."
    except ValueError:
        return "."


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 4 — EMBEDDING  (OpenAI text-embedding-3-small)
# ═════════════════════════════════════════════════════════════════════════════

def embed_texts(texts: List[str], retries: int = 3) -> List[List[float]]:
    all_embeddings: List[List[float]] = []

    for i in range(0, len(texts), EMBEDDING_BATCH):
        batch = [t.replace("\n", " ").strip() for t in texts[i: i + EMBEDDING_BATCH]]

        for attempt in range(retries):
            try:
                resp = openai_client().embeddings.create(
                    model=EMBEDDING_MODEL,
                    input=batch,
                )
                all_embeddings.extend(item.embedding for item in resp.data)
                break
            except Exception as e:
                if attempt == retries - 1:
                    raise
                wait = 2 ** attempt
                print(f"    [WARN] Embedding batch failed ({e}). Retry in {wait}s…")
                time.sleep(wait)

    return all_embeddings


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 5 — QDRANT  (collection setup + upsert)
# ═════════════════════════════════════════════════════════════════════════════

def setup_collection(recreate: bool = False) -> None:
    client = qdrant_client()

    if recreate and client.collection_exists(COLLECTION_NAME):
        client.delete_collection(COLLECTION_NAME)
        print(f"  Deleted existing collection '{COLLECTION_NAME}'")

    if not client.collection_exists(COLLECTION_NAME):
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(size=EMBEDDING_DIM, distance=Distance.COSINE),
        )
        print(f"  Created collection '{COLLECTION_NAME}' ({EMBEDDING_DIM}-dim, cosine)")

        keyword_fields = [
            "doc_type", "industry_vertical", "function_area",
            "solution_type", "section_type", "extraction_method", "source_folder",
        ]
        for field in keyword_fields:
            client.create_payload_index(COLLECTION_NAME, field, PayloadSchemaType.KEYWORD)
        client.create_payload_index(COLLECTION_NAME, "has_roi_metrics", PayloadSchemaType.BOOL)
        print(f"  Payload indexes created: {keyword_fields + ['has_roi_metrics']}")
    else:
        print(f"  Collection '{COLLECTION_NAME}' already exists — skipping setup")


def upsert_chunks(chunks: List[Dict], embeddings: List[List[float]]) -> None:
    points = [
        PointStruct(
            id      = str(uuid.uuid4()),
            vector  = vector,
            payload = {k: v for k, v in chunk.items()},
        )
        for chunk, vector in zip(chunks, embeddings)
    ]
    qdrant_client().upsert(collection_name=COLLECTION_NAME, points=points)


# ═════════════════════════════════════════════════════════════════════════════
# SECTION 6 — PIPELINE ORCHESTRATION
# ═════════════════════════════════════════════════════════════════════════════

def _should_skip(path: Path, processed_stems: Set[str]) -> tuple[bool, str]:
    name = path.name
    ext  = path.suffix.lower()

    if any(skip in path.parts for skip in SKIP_FOLDERS):
        return True, "excluded folder"
    if ext in SKIP_EXTENSIONS:
        return True, "unsupported extension"
    if name in SKIP_DUPLICATES:
        return True, "known duplicate"
    if ext == ".pdf" and name in SKIP_PDF_USE_PPTX:
        return True, "image-heavy PDF — PPTX twin preferred"
    if ext == ".pptx" and path.stem.lower() in processed_stems:
        return True, "PDF twin already processed"

    return False, ""


def process_file(path: Path) -> List[Dict]:
    """Extract → chunk → enrich one file. Returns enriched chunk list."""
    pages = extract_file(path)
    if not pages:
        return []

    # Determine doc_type early (needed for chunk strategy)
    doc_type = "capability_deck"
    folder_parts = [p.lower() for p in path.parts]
    for folder_key, dtype in FOLDER_TO_DOCTYPE.items():
        if folder_key.lower() in " ".join(folder_parts):
            doc_type = dtype
            break

    chunks = chunk_pages(pages, doc_type)
    if not chunks:
        return []

    doc_first500 = " ".join(p["text"] for p in pages)[:500]
    rel_path     = _relative_path(path)
    doc_id       = hashlib.md5(rel_path.encode()).hexdigest()[:12]

    enriched = []
    for idx, chunk in enumerate(chunks):
        meta              = enrich_chunk(path, chunk, doc_first500)
        meta["doc_id"]    = doc_id
        meta["chunk_id"]  = f"{doc_id}_c{idx:03d}"
        meta["chunk_index"]  = idx
        meta["total_chunks"] = len(chunks)
        enriched.append(meta)

    vlm_count = sum(1 for c in enriched if c.get("extraction_method") == "vlm")
    print(
        f"    → {len(enriched)} chunks | "
        f"doc_type={enriched[0]['doc_type']} | "
        f"verticals={enriched[0]['industry_vertical']}"
        + (f" | vlm_pages={vlm_count}" if vlm_count else "")
    )
    return enriched


def run(data_dir: Path, recreate: bool = False, dry_run: bool = False, single_file: str | None = None) -> None:
    print(f"\n{'=' * 64}")
    print("MindMap Sales Collateral — RAG Ingestion Pipeline")
    print(f"Data dir  : {data_dir}")
    print(f"Dry run   : {dry_run}")
    print(f"{'=' * 64}\n")

    if not dry_run:
        setup_collection(recreate=recreate)

    # ── File discovery ─────────────────────────────────────────────────────
    if single_file:
        all_files = [data_dir / single_file] if not Path(single_file).is_absolute() else [Path(single_file)]
    else:
        all_files = sorted(f for f in data_dir.rglob("*") if f.is_file())

    # ── Process files ──────────────────────────────────────────────────────
    processed_stems: Set[str] = set()
    all_chunks: List[Dict]    = []
    stats = {"processed": 0, "skipped": 0, "vlm_files": 0}

    for path in all_files:
        skip, reason = _should_skip(path, processed_stems)
        if skip:
            print(f"  [SKIP] {path.name:<60} {reason}")
            stats["skipped"] += 1
            continue

        is_vlm = path.name in VLM_REQUIRED_FILES
        print(f"\n  [{'VLM' if is_vlm else path.suffix[1:].upper():>4}] {path.name}")
        if is_vlm:
            stats["vlm_files"] += 1

        chunks = process_file(path)
        if not chunks:
            continue

        # Track stems so PPTX twins of text-PDFs are skipped
        if path.suffix.lower() == ".pdf" and path.name not in SKIP_PDF_USE_PPTX:
            processed_stems.add(path.stem.lower())

        all_chunks.extend(chunks)
        stats["processed"] += 1

    # ── Summary ────────────────────────────────────────────────────────────
    print(f"\n{'=' * 64}")
    print(f"Extraction complete")
    print(f"  Files processed : {stats['processed']}")
    print(f"  Files skipped   : {stats['skipped']}")
    print(f"  VLM files       : {stats['vlm_files']}")
    print(f"  Total chunks    : {len(all_chunks)}")

    if dry_run:
        preview_path = data_dir.parent / "chunks_preview.json"
        with open(preview_path, "w", encoding="utf-8") as f:
            json.dump(all_chunks[:20], f, indent=2, ensure_ascii=False)
        print(f"\n[DRY RUN] Preview saved → {preview_path}")
        return

    # ── Embed + upsert ─────────────────────────────────────────────────────
    print(f"\nEmbedding + uploading to Qdrant …")
    t0          = time.time()
    total_batches = (len(all_chunks) - 1) // UPSERT_BATCH + 1

    for i in range(0, len(all_chunks), UPSERT_BATCH):
        batch     = all_chunks[i: i + UPSERT_BATCH]
        batch_num = i // UPSERT_BATCH + 1
        print(f"  Batch {batch_num}/{total_batches} — {len(batch)} chunks … ", end="", flush=True)
        embeddings = embed_texts([c["text"] for c in batch])
        upsert_chunks(batch, embeddings)
        print("done")

    info = qdrant_client().get_collection(COLLECTION_NAME)
    print(f"\n{'=' * 64}")
    print(f"Ingestion complete in {time.time() - t0:.1f}s")
    print(f"  Qdrant vectors  : {info.vectors_count}")
    print(f"{'=' * 64}\n")


# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="MindMap RAG ingestion pipeline")
    parser.add_argument("--data-dir",  default=str(DATA_DIR),  help="Path to Sales Collateral folder")
    parser.add_argument("--recreate",  action="store_true",     help="Drop and recreate Qdrant collection")
    parser.add_argument("--dry-run",   action="store_true",     help="Extract + chunk only, skip embedding/Qdrant")
    parser.add_argument("--file",      default=None,            help="Process a single file (relative to data-dir)")
    args = parser.parse_args()

    run(
        data_dir    = Path(args.data_dir),
        recreate    = args.recreate,
        dry_run     = args.dry_run,
        single_file = args.file,
    )
